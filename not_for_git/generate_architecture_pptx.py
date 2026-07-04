# Databricks notebook source
# /// script
# [tool.databricks.environment]
# environment_version = "5"
# ///
# DBTITLE 1,Success Message
print("="*80)
print("\n✅ SUCCESS! PowerPoint presentation generated successfully!\n")
print("File location:")
print(f"  {output_path}")
print("\nPresentation details:")
print(f"  • Total slides: {len(prs.slides)}")
print(f"  • Slide dimensions: {prs.slide_width.inches}\" x {prs.slide_height.inches}\"")
print("  • Background: Dark blue (#1a2332)")
print("  • Color scheme: Teal (Free) / Orange (Databricks)")
print("  • All shapes are editable native PowerPoint objects")
print("\nSlides included:")
print("  1. Title Slide")
print("  2. Complete 4-Layer Architecture Overview")
print("  3. Layers 1 & 2 Detailed (Document Ingestion + RAG Pipeline)")
print("  4. Layer 3 - Multi-Agent Orchestration (LangGraph)")
print("  5. Layer 4 & Key Architectural Decisions")
print("\n" + "="*80)

# COMMAND ----------

# DBTITLE 1,Save Presentation
# Save the presentation
output_path = '/Workspace/Users/gb.burcea@gmail.com/agentic_quality_check/pptx/agentic_quality_check_architecture.pptx'
prs.save(output_path)
print(f"Presentation saved to: {output_path}")

# COMMAND ----------

# DBTITLE 1,Create Presentation and All Slides
# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
slide1_layout = prs.slide_layouts[6]  # Blank layout
slide1 = prs.slides.add_slide(slide1_layout)
set_background(slide1, COLORS['background'])

# Main title
add_title_shape(slide1, "Agentic RAG Quality Checker Architecture", Inches(2.5), font_size=Pt(40))

# Subtitle 1
add_title_shape(slide1, "LangChain + LangGraph Multi-Agent System", Inches(3.5), font_size=Pt(24))

# Subtitle 2 with color-coded parts
subtitle2_left = Inches(1.5)
subtitle2_top = Inches(4.5)
subtitle2_width = Inches(7)

txBox = slide1.shapes.add_textbox(subtitle2_left, subtitle2_top, subtitle2_width, Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER

# Add "Dual-Path Design: "
run1 = p.add_run()
run1.text = "Dual-Path Design: "
run1.font.size = Pt(20)
run1.font.color.rgb = COLORS['white']
run1.font.name = 'Segoe UI'

# Add "Free/Open-Source" in teal
run2 = p.add_run()
run2.text = "Free/Open-Source"
run2.font.size = Pt(20)
run2.font.color.rgb = COLORS['free']
run2.font.bold = True
run2.font.name = 'Segoe UI'

# Add " ↔ "
run3 = p.add_run()
run3.text = " ↔ "
run3.font.size = Pt(20)
run3.font.color.rgb = COLORS['white']
run3.font.name = 'Segoe UI'

# Add "Databricks Enterprise" in orange
run4 = p.add_run()
run4.text = "Databricks Enterprise"
run4.font.size = Pt(20)
run4.font.color.rgb = COLORS['databricks']
run4.font.bold = True
run4.font.name = 'Segoe UI'

print("✓ Slide 1: Title Slide created")

# ============================================================================
# SLIDE 2: Complete 4-Layer Architecture Overview
# ============================================================================
slide2_layout = prs.slide_layouts[6]
slide2 = prs.slides.add_slide(slide2_layout)
set_background(slide2, COLORS['background'])

# Title
add_title_shape(slide2, "Complete 4-Layer Architecture Overview", Inches(0.3), font_size=Pt(28))

# Layer positioning
layer_left = Inches(1.5)
layer_width = Inches(7)
layer_height = Inches(1.2)
layer_start_top = Inches(1.3)
layer_spacing = Inches(0.3)

layers = [
    {
        'title': 'Layer 1: Document Ingestion & Preprocessing',
        'free': 'PyPDF2\nLocal Storage',
        'db': 'LangChain Loader\nUC Volumes'
    },
    {
        'title': 'Layer 2: RAG Pipeline',
        'free': 'sentence-transformers\nFAISS',
        'db': 'Foundation Models\nVector Search'
    },
    {
        'title': 'Layer 3: Multi-Agent Orchestration',
        'free': 'HuggingFace Models\nLocal Execution',
        'db': 'Foundation Models API\nServerless'
    },
    {
        'title': 'Layer 4: User Interface',
        'free': 'Streamlit\nLocal/Cloud',
        'db': 'Databricks Apps\nWorkspace'
    }
]

for i, layer in enumerate(layers):
    current_top = layer_start_top + i * (layer_height + layer_spacing)
    
    # Layer title background
    add_rounded_rect(slide2, layer_left, current_top, layer_width, Inches(0.35), 
                    layer['title'], COLORS['light_gray'], COLORS['background'], Pt(12))
    
    # Split box: Free (left) and Databricks (right)
    box_top = current_top + Inches(0.4)
    box_height = Inches(0.75)
    box_width = Inches(3.3)
    
    # Free option (left, teal)
    add_rounded_rect(slide2, layer_left + Inches(0.2), box_top, box_width, box_height,
                    layer['free'], COLORS['free'], COLORS['background'], Pt(11))
    
    # Databricks option (right, orange)
    add_rounded_rect(slide2, layer_left + Inches(3.7), box_top, box_width, box_height,
                    layer['db'], COLORS['databricks'], COLORS['background'], Pt(11))
    
    # Arrow to next layer (if not last)
    if i < len(layers) - 1:
        arrow_x = layer_left + layer_width / 2
        arrow_y1 = current_top + layer_height
        arrow_y2 = current_top + layer_height + layer_spacing
        add_arrow(slide2, arrow_x, arrow_y1, arrow_x, arrow_y2)

print("✓ Slide 2: 4-Layer Architecture Overview created")

# ============================================================================
# SLIDE 3: Layer 1 & 2 Detailed
# ============================================================================
slide3_layout = prs.slide_layouts[6]
slide3 = prs.slides.add_slide(slide3_layout)
set_background(slide3, COLORS['background'])

# Title
add_title_shape(slide3, "Layers 1 & 2: Document Ingestion + RAG Pipeline", Inches(0.3), font_size=Pt(26))

# Layer 1 Section
add_text_shape(slide3, "Layer 1: Document Ingestion & Preprocessing", Inches(1.2), Inches(0.5),
              Inches(9), Inches(0.4), Pt(18), PP_ALIGN.LEFT, True)

# PDF flow
flow_top = Inches(1.7)
flow_left = Inches(1)
box_width = Inches(1.5)
box_height = Inches(0.6)
spacing = Inches(0.3)

# PDF Files
add_rounded_rect(slide3, flow_left, flow_top, box_width, box_height, "PDF Files", 
                COLORS['light_gray'], COLORS['background'], Pt(12))
add_arrow(slide3, flow_left + box_width, flow_top + box_height/2,
         flow_left + box_width + spacing, flow_top + box_height/2)

# PyPDF2 Loader
add_rounded_rect(slide3, flow_left + box_width + spacing, flow_top, box_width, box_height,
                "PyPDF2\nLoader", COLORS['free'], COLORS['background'], Pt(12))
add_arrow(slide3, flow_left + 2*(box_width + spacing), flow_top + box_height/2,
         flow_left + 2*(box_width + spacing) + spacing, flow_top + box_height/2)

# Text Splitter
add_rounded_rect(slide3, flow_left + 2*(box_width + spacing) + spacing, flow_top, box_width, box_height,
                "Text\nSplitter", COLORS['light_gray'], COLORS['background'], Pt(12))
add_arrow(slide3, flow_left + 3*(box_width + spacing) + spacing, flow_top + box_height/2,
         flow_left + 3*(box_width + spacing) + 2*spacing, flow_top + box_height/2)

# Chunks
add_rounded_rect(slide3, flow_left + 3*(box_width + spacing) + 2*spacing, flow_top, box_width, box_height,
                "Text\nChunks", COLORS['light_gray'], COLORS['background'], Pt(12))

# Storage options
storage_top = Inches(2.5)
add_text_shape(slide3, "Storage:", storage_top, Inches(0.5), Inches(2), Inches(0.3), Pt(12), PP_ALIGN.LEFT, True)
add_text_shape(slide3, "• Free: Local filesystem", storage_top + Inches(0.25), Inches(0.7), Inches(4), Inches(0.3), Pt(11), color=COLORS['free'])
add_text_shape(slide3, "• Databricks: Unity Catalog Volumes", storage_top + Inches(0.25), Inches(5), Inches(4.5), Inches(0.3), Pt(11), color=COLORS['databricks'])

# Layer 2 Section
add_text_shape(slide3, "Layer 2: RAG Pipeline (Embedding + Vector Store + Retrieval)", Inches(3.4), Inches(0.5),
              Inches(9), Inches(0.4), Pt(18), PP_ALIGN.LEFT, True)

# Three components side by side
comp_top = Inches(3.9)
comp_left = Inches(0.7)
comp_width = Inches(2.7)
comp_height = Inches(2.2)
comp_spacing = Inches(0.3)

# Embedding
add_rounded_rect(slide3, comp_left, comp_top, comp_width, Inches(0.4), "Embedding Generation",
                COLORS['light_gray'], COLORS['background'], Pt(12))
add_rounded_rect(slide3, comp_left + Inches(0.1), comp_top + Inches(0.5), comp_width - Inches(0.2), Inches(0.7),
                "Free:\nsentence-transformers\n384-dim vectors", COLORS['free'], COLORS['background'], Pt(10))
add_rounded_rect(slide3, comp_left + Inches(0.1), comp_top + Inches(1.3), comp_width - Inches(0.2), Inches(0.7),
                "Databricks:\nFoundation Models\n1024-dim vectors", COLORS['databricks'], COLORS['background'], Pt(10))

# Vector Store
add_rounded_rect(slide3, comp_left + comp_width + comp_spacing, comp_top, comp_width, Inches(0.4),
                "Vector Indexing", COLORS['light_gray'], COLORS['background'], Pt(12))
add_rounded_rect(slide3, comp_left + comp_width + comp_spacing + Inches(0.1), comp_top + Inches(0.5),
                comp_width - Inches(0.2), Inches(0.7), "Free:\nFAISS\nIn-memory, fast",
                COLORS['free'], COLORS['background'], Pt(10))
add_rounded_rect(slide3, comp_left + comp_width + comp_spacing + Inches(0.1), comp_top + Inches(1.3),
                comp_width - Inches(0.2), Inches(0.7), "Databricks:\nVector Search\nPersistent, scalable",
                COLORS['databricks'], COLORS['background'], Pt(10))

# Retrieval
add_rounded_rect(slide3, comp_left + 2*(comp_width + comp_spacing), comp_top, comp_width, Inches(0.4),
                "Retrieval", COLORS['light_gray'], COLORS['background'], Pt(12))
add_rounded_rect(slide3, comp_left + 2*(comp_width + comp_spacing) + Inches(0.1), comp_top + Inches(0.5),
                comp_width - Inches(0.2), Inches(1.5), "Query Embedding\n↓\nSimilarity Search\n↓\nTop-K Documents\n(with metadata)",
                COLORS['light_gray'], COLORS['background'], Pt(10))

print("✓ Slide 3: Layers 1 & 2 Detailed created")

# ============================================================================
# SLIDE 4: Layer 3 - Multi-Agent Orchestration
# ============================================================================
slide4_layout = prs.slide_layouts[6]
slide4 = prs.slides.add_slide(slide4_layout)
set_background(slide4, COLORS['background'])

# Title
add_title_shape(slide4, "Layer 3: Multi-Agent Orchestration (LangGraph)", Inches(0.3), font_size=Pt(26))

# Supervisor Agent (top center)
sup_left = Inches(3.5)
sup_top = Inches(1.5)
sup_width = Inches(3)
sup_height = Inches(0.7)
add_rounded_rect(slide4, sup_left, sup_top, sup_width, sup_height, "SUPERVISOR AGENT\n(Routes to workers)",
                COLORS['light_gray'], COLORS['background'], Pt(14))

# Worker agents (bottom row)
worker_top = Inches(3.5)
worker_width = Inches(2.5)
worker_height = Inches(1.2)
worker_spacing = Inches(0.5)
worker_start = Inches(1.3)

# Retrieval Agent
add_rounded_rect(slide4, worker_start, worker_top, worker_width, worker_height,
                "RETRIEVAL\nAGENT\n\nCalls Layer 2\nRAG Pipeline", COLORS['free'], COLORS['background'], Pt(12))

# Response Generator Agent
add_rounded_rect(slide4, worker_start + worker_width + worker_spacing, worker_top, worker_width, worker_height,
                "RESPONSE\nGENERATOR\n\nContext + Question\n→ Answer", COLORS['databricks'], COLORS['background'], Pt(12))

# Quality Checker Agent
add_rounded_rect(slide4, worker_start + 2*(worker_width + worker_spacing), worker_top, worker_width, worker_height,
                "QUALITY\nCHECKER\n\nEvaluates:\nAccuracy, Citations", COLORS['free'], COLORS['background'], Pt(12))

# Arrows from Supervisor to workers
sup_center_x = sup_left + sup_width / 2
sup_bottom_y = sup_top + sup_height

# To Retrieval Agent
add_arrow(slide4, sup_center_x - Inches(0.8), sup_bottom_y, worker_start + worker_width/2, worker_top)

# To Response Generator
add_arrow(slide4, sup_center_x, sup_bottom_y, worker_start + worker_width + worker_spacing + worker_width/2, worker_top)

# To Quality Checker
add_arrow(slide4, sup_center_x + Inches(0.8), sup_bottom_y, worker_start + 2*(worker_width + worker_spacing) + worker_width/2, worker_top)

# Feedback loop arrow (Quality Checker → Response Generator)
feedback_x1 = worker_start + 2*(worker_width + worker_spacing) + worker_width
feedback_y = worker_top + worker_height/2
feedback_x2 = worker_start + worker_width + worker_spacing + worker_width
add_arrow(slide4, feedback_x1, feedback_y, feedback_x2, feedback_y)

# Feedback label
add_text_shape(slide4, "Feedback loop\n(if fail)", Inches(5.5), worker_start + 2*(worker_width + worker_spacing) + worker_width - Inches(0.3),
              Inches(1.5), Inches(0.6), Pt(9), PP_ALIGN.CENTER, color=COLORS['arrow'])

# LLM Options
llm_top = Inches(5.5)
add_text_shape(slide4, "LLM Options:", llm_top, Inches(0.5), Inches(2), Inches(0.3), Pt(14), PP_ALIGN.LEFT, True)
add_text_shape(slide4, "• Free: HuggingFace (Mistral-7B, Llama-2) - Local inference",
              llm_top + Inches(0.35), Inches(0.7), Inches(9), Inches(0.3), Pt(11), color=COLORS['free'])
add_text_shape(slide4, "• Databricks: Foundation Models API (DBRX, Llama 3.1, Mixtral) - Serverless",
              llm_top + Inches(0.65), Inches(0.7), Inches(9), Inches(0.3), Pt(11), color=COLORS['databricks'])

# State Management
state_top = Inches(6.4)
add_text_shape(slide4, "State Management (LangGraph): Conversation history • Retrieved docs • Quality scores • Loop counter",
              state_top, Inches(0.5), Inches(9), Inches(0.4), Pt(10), PP_ALIGN.CENTER, color=COLORS['light_gray'])

print("✓ Slide 4: Layer 3 - Multi-Agent Orchestration created")

# ============================================================================
# SLIDE 5: Layer 4 & Key Decisions
# ============================================================================
slide5_layout = prs.slide_layouts[6]
slide5 = prs.slides.add_slide(slide5_layout)
set_background(slide5, COLORS['background'])

# Title
add_title_shape(slide5, "Layer 4 & Key Architectural Decisions", Inches(0.3), font_size=Pt(26))

# Layer 4 Section
add_text_shape(slide5, "Layer 4: User Interface", Inches(1.2), Inches(0.5), Inches(9), Inches(0.4),
              Pt(18), PP_ALIGN.LEFT, True)

ui_top = Inches(1.7)
ui_left = Inches(1.5)
ui_width = Inches(3.3)
ui_height = Inches(1.5)

# Streamlit (Free)
add_rounded_rect(slide5, ui_left, ui_top, ui_width, Inches(0.4), "Streamlit (Free)",
                COLORS['free'], COLORS['background'], Pt(14))
add_text_shape(slide5, "• Chat interface\n• Document upload\n• Agent trace\n• Local or Streamlit Cloud",
              ui_top + Inches(0.5), ui_left + Inches(0.2), ui_width - Inches(0.4), Inches(1),
              Pt(11), PP_ALIGN.LEFT, color=COLORS['white'])

# Databricks Apps
add_rounded_rect(slide5, ui_left + ui_width + Inches(0.6), ui_top, ui_width, Inches(0.4),
                "Databricks Apps", COLORS['databricks'], COLORS['background'], Pt(14))
add_text_shape(slide5, "• Workspace integration\n• Unity Catalog auth\n• Direct data access\n• Enterprise deployment",
              ui_top + Inches(0.5), ui_left + ui_width + Inches(0.6) + Inches(0.2),
              ui_width - Inches(0.4), Inches(1), Pt(11), PP_ALIGN.LEFT, color=COLORS['white'])

# Key Architectural Decisions
add_text_shape(slide5, "Key Architectural Decisions", Inches(3.5), Inches(0.5), Inches(9), Inches(0.4),
              Pt(18), PP_ALIGN.LEFT, True)

decisions_top = Inches(4)
decisions_left = Inches(0.8)
decision_width = Inches(4.2)
decision_height = Inches(1.2)
decision_spacing = Inches(0.3)

decisions = [
    {
        'title': '1. Dual-Path Design',
        'desc': 'Every component has free and\nDatabricks options →\nDevelop on free, scale with DB'
    },
    {
        'title': '2. LangChain Abstraction',
        'desc': 'Unified API regardless of backend\n→ Switch FAISS to Vector Search\nwithout code changes'
    },
    {
        'title': '3. LangGraph Orchestration',
        'desc': 'Stateful, cyclic workflows →\nQuality feedback loops impossible\nwith simple chains'
    },
    {
        'title': '4. Modular Layers',
        'desc': 'Clear input/output contracts →\nEasy to test, debug, and\noptimize independently'
    }
]

for i, decision in enumerate(decisions):
    row = i // 2
    col = i % 2
    
    left = decisions_left + col * (decision_width + decision_spacing)
    top = decisions_top + row * (decision_height + decision_spacing)
    
    # Title bar
    add_rounded_rect(slide5, left, top, decision_width, Inches(0.35), decision['title'],
                    COLORS['databricks'], COLORS['white'], Pt(12))
    
    # Description box
    add_text_shape(slide5, decision['desc'], top + Inches(0.4), left + Inches(0.2),
                  decision_width - Inches(0.4), Inches(0.75), Pt(10), PP_ALIGN.LEFT, color=COLORS['white'])

print("✓ Slide 5: Layer 4 & Key Decisions created")
print("\n✓ All 5 slides created successfully!")

# COMMAND ----------

# DBTITLE 1,Install python-pptx
# MAGIC %pip install python-pptx

# COMMAND ----------

# DBTITLE 1,Imports and Helper Functions
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Color Scheme
COLORS = {
    'background': RGBColor(26, 35, 50),      # Dark blue #1a2332
    'free': RGBColor(0, 212, 170),            # Teal #00d4aa
    'databricks': RGBColor(255, 107, 53),     # Orange #ff6b35
    'white': RGBColor(255, 255, 255),
    'light_gray': RGBColor(200, 200, 200),
    'arrow': RGBColor(150, 150, 150)
}

def set_background(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_shape(slide, text, top, left=Inches(0.5), width=Inches(9), height=Inches(0.8), font_size=Pt(32)):
    """Add a title text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = font_size
    p.font.color.rgb = COLORS['white']
    p.font.bold = True
    p.font.name = 'Segoe UI'
    return txBox

def add_text_shape(slide, text, top, left, width, height, font_size=Pt(14), align=PP_ALIGN.LEFT, bold=False, color=None):
    """Add a text box with custom styling"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.font.size = font_size
    p.font.color.rgb = color if color else COLORS['white']
    p.font.bold = bold
    p.font.name = 'Segoe UI'
    return txBox

def add_rounded_rect(slide, left, top, width, height, text, fill_color, text_color=None, font_size=Pt(14)):
    """Add a rounded rectangle with text"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    
    # Set fill color
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    
    # Set border
    shape.line.color.rgb = COLORS['white']
    shape.line.width = Pt(1.5)
    
    # Add text
    if text:
        tf = shape.text_frame
        tf.text = text
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = font_size
        p.font.color.rgb = text_color if text_color else COLORS['white']
        p.font.bold = True
        p.font.name = 'Segoe UI'
    
    return shape

def add_arrow(slide, x1, y1, x2, y2):
    """Add a connector arrow between two points"""
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        x1, y1, x2, y2
    )
    connector.line.color.rgb = COLORS['arrow']
    connector.line.width = Pt(2)
    return connector

print("✓ Imports and helper functions loaded")
